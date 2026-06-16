"""
Text extraction helper for PT Study Agent.
Usage: python extract_text.py --file <path> --output <path>
Outputs JSON: { success, filename, file_type, page_count, word_count, text,
                scanned, image_paths, pages_dir, capped, images, images_skipped,
                error }
  images        — for .html: [{path, alt}] extracted figures saved under pages_dir
  images_skipped — count of <img> skipped (external URLs, svg, unreadable)
"""
import argparse
import base64
import html as htmllib
import json
import os
import pathlib
import re
import sys
from html.parser import HTMLParser

from lkconfig import get as cfg

SCRIPTS_DIR = pathlib.Path(__file__).parent
SCANNED_WORDS_PER_PAGE_THRESHOLD = cfg("scanned_words_per_page_threshold")
MAX_SCANNED_PAGES = cfg("max_scanned_pages")  # default render cap; override with --max-pages


def _safe_name(name: str) -> str:
    """Sanitize a filename stem for use as a path component.

    Windows forbids the chars <>:"/\\|?* and silently strips trailing
    spaces/dots from directory names, which breaks rasterizers that write
    to the un-stripped path (e.g. 'The Bones of the Foot ' → page write fails).
    """
    name = name.strip()
    for ch in '<>:"/\\|?*':
        name = name.replace(ch, "_")
    name = name.rstrip(" .")
    return name or "doc"


def extract_pdf(path):
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        total_pages = len(pdf.pages)
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    full_text = "\n\n".join(pages)
    word_count = len(full_text.split())
    is_scanned = total_pages > 0 and word_count < SCANNED_WORDS_PER_PAGE_THRESHOLD * total_pages
    return full_text, total_pages, is_scanned


def render_pdf_pages(path, max_pages=MAX_SCANNED_PAGES):
    import fitz
    basename = _safe_name(pathlib.Path(path).stem)
    out_dir = SCRIPTS_DIR / "tmp_pages" / basename
    out_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(path)
    total_pages = len(doc)
    cap = total_pages if max_pages <= 0 else min(total_pages, max_pages)
    capped = total_pages > cap
    render_count = cap

    image_paths = []
    scale = cfg("render_scale")
    mat = fitz.Matrix(scale, scale)
    for i in range(render_count):
        pix = doc[i].get_pixmap(matrix=mat)
        img_path = out_dir / f"page_{i + 1:03d}.png"
        pix.save(str(img_path))
        image_paths.append(str(img_path))

    doc.close()
    return image_paths, total_pages, capped, str(out_dir)


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides), len(prs.slides)


def extract_docx(path):
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs), len(paragraphs)


def extract_text_file(path):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return text, text.count("\n") + 1


_DATA_URI_RE = re.compile(
    r"^data:image/(?P<fmt>[\w.+-]+);base64,(?P<data>.+)$", re.IGNORECASE | re.DOTALL)
_RASTER_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}


def _img_ext_from_fmt(fmt):
    return {"jpeg": "jpg", "svg+xml": "svg"}.get(fmt.lower(), fmt.lower())


class _HTMLExtractor(HTMLParser):
    """Pull readable text + <img> (src, alt) from HTML; drop script/style/head."""
    _SKIP = {"script", "style", "head", "noscript", "svg"}
    _BLOCK = {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6",
              "section", "article", "table", "ul", "ol", "header", "footer",
              "blockquote", "pre"}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self._skip_depth = 0
        self.text_parts = []
        self.images = []  # [(src, alt)] in document order

    def handle_starttag(self, tag, attrs):
        if tag in self._SKIP:
            self._skip_depth += 1
            return
        if tag == "img":
            d = dict(attrs)
            src = (d.get("src") or "").strip()
            if src:
                self.images.append((src, (d.get("alt") or "").strip()))
        if tag in self._BLOCK:
            self.text_parts.append("\n")

    def handle_startendtag(self, tag, attrs):
        # self-closing <img .../> still yields an image
        self.handle_starttag(tag, attrs)

    def handle_endtag(self, tag):
        if tag in self._SKIP and self._skip_depth > 0:
            self._skip_depth -= 1
        elif tag in self._BLOCK:
            self.text_parts.append("\n")

    def handle_data(self, data):
        if self._skip_depth == 0:
            self.text_parts.append(data)

    def get_text(self):
        lines = [" ".join(ln.split()) for ln in "".join(self.text_parts).splitlines()]
        return "\n".join(ln for ln in lines if ln)


def extract_html(path):
    """Return (text, images, images_skipped, pages_dir).

    images = [{path, alt}] for each embedded (base64) or local-file <img> saved
    under SCRIPTS_DIR/tmp_html/<stem>. External (http/protocol-relative) images
    and SVG/unreadable ones are skipped (counted). pages_dir is None if nothing
    was saved (so the caller has nothing to clean up)."""
    import shutil
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        raw = f.read()
    parser = _HTMLExtractor()
    parser.feed(raw)
    text = parser.get_text()

    stem = _safe_name(pathlib.Path(path).stem)
    out_dir = SCRIPTS_DIR / "tmp_html" / stem
    base_dir = pathlib.Path(path).resolve().parent

    images, skipped, saved = [], 0, 0
    for src, alt in parser.images:
        src_u = htmllib.unescape(src)
        try:
            m = _DATA_URI_RE.match(src_u)
            if m:
                fmt = _img_ext_from_fmt(m.group("fmt"))
                if fmt == "svg":                      # PIL can't rasterize SVG
                    skipped += 1
                    continue
                blob = base64.b64decode(m.group("data"), validate=False)
                out_dir.mkdir(parents=True, exist_ok=True)
                dest = out_dir / f"img_{saved + 1:03d}.{fmt}"
                dest.write_bytes(blob)
                saved += 1
                images.append({"path": str(dest), "alt": htmllib.unescape(alt)})
            elif re.match(r"^(https?:)?//", src_u, re.IGNORECASE) or src_u.lower().startswith("data:"):
                skipped += 1                            # remote URL or non-image data URI
            else:
                local = (base_dir / src_u).resolve()
                if local.is_file() and local.suffix.lower() in _RASTER_EXT:
                    out_dir.mkdir(parents=True, exist_ok=True)
                    dest = out_dir / f"img_{saved + 1:03d}{local.suffix.lower()}"
                    shutil.copyfile(local, dest)
                    saved += 1
                    images.append({"path": str(dest), "alt": htmllib.unescape(alt)})
                else:
                    skipped += 1
        except Exception:
            skipped += 1
    return text, images, skipped, (str(out_dir) if saved else None)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--max-pages", type=int, default=MAX_SCANNED_PAGES,
                        help="scanned-PDF page render cap; 0 = no cap")
    args = parser.parse_args()

    result = {
        "success": False,
        "filename": os.path.basename(args.file),
        "file_type": None,
        "page_count": 0,
        "word_count": 0,
        "text": "",
        "scanned": False,
        "image_paths": [],
        "pages_dir": None,
        "capped": False,
        "images": [],
        "images_skipped": 0,
        "error": None,
    }

    ext = os.path.splitext(args.file)[1].lower()

    try:
        if not os.path.exists(args.file):
            raise FileNotFoundError(f"File not found: {args.file}")

        if ext == ".pdf":
            result["file_type"] = "pdf"
            text, page_count, is_scanned = extract_pdf(args.file)
            result["page_count"] = page_count
            if is_scanned:
                result["scanned"] = True
                image_paths, _, capped, pages_dir = render_pdf_pages(args.file, args.max_pages)
                result["image_paths"] = image_paths
                result["pages_dir"] = pages_dir
                result["capped"] = capped
            else:
                result["text"] = text
                result["word_count"] = len(text.split())
        elif ext == ".pptx":
            result["file_type"] = "pptx"
            result["text"], result["page_count"] = extract_pptx(args.file)
        elif ext == ".docx":
            result["file_type"] = "docx"
            result["text"], result["page_count"] = extract_docx(args.file)
        elif ext in (".txt", ".md"):
            result["file_type"] = ext.lstrip(".")
            result["text"], result["page_count"] = extract_text_file(args.file)
        elif ext in (".html", ".htm"):
            result["file_type"] = "html"
            text, images, skipped, pages_dir = extract_html(args.file)
            result["text"] = text
            result["images"] = images
            result["images_skipped"] = skipped
            result["pages_dir"] = pages_dir
            result["page_count"] = 1
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        if not result["scanned"]:
            result["word_count"] = len(result["text"].split())
        result["success"] = True

    except Exception as e:
        result["error"] = str(e)

    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)


if __name__ == "__main__":
    main()
